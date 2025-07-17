"""
This utility module contains functions to generate standardized reports.

It now includes a function for a single-project status report and a more
comprehensive, board-ready, multi-slide portfolio summary deck, fulfilling a key
requirement of the enhanced Communication & Reporting Toolkit.
"""
import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
# --- FIX: Corrected the import from MSO_ALIGN to PP_ALIGN ---
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from datetime import date
import plotly.io as pio
import plotly.express as px

# --- Helper Functions ---

def set_cell_text(cell, text, bold=False, font_size=10, align='LEFT'):
    """Helper function to set text and formatting in a table cell."""
    cell.text_frame.clear()
    p = cell.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    
    # --- FIX: Updated the logic to use the correct PP_ALIGN enumeration ---
    alignment_map = {
        'LEFT': PP_ALIGN.LEFT,
        'CENTER': PP_ALIGN.CENTER,
        'RIGHT': PP_ALIGN.RIGHT,
        'JUSTIFY': PP_ALIGN.JUSTIFY
    }
    p.alignment = alignment_map.get(align.upper(), PP_ALIGN.LEFT)
    
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = RGBColor(0, 0, 0)
    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

def add_plotly_fig_to_slide(fig, slide, left, top, width, height):
    """Helper to save a Plotly figure as an image and add it to a slide."""
    img_bytes = pio.to_image(fig, format="png", scale=2)
    return slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width, height)

# --- Report Generation Functions ---

def generate_project_status_report(project_details: dict, milestones_df: pd.DataFrame, risks_df: pd.DataFrame) -> io.BytesIO:
    """Generates a detailed, one-page PowerPoint status report for a given project."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(16), Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout

    # --- Title and Date ---
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.75))
    p = title_shape.text_frame.paragraphs[0]
    p.text = f"Project Status Report: {project_details.get('name', 'N/A')}"
    p.font.size, p.font.bold = Pt(32), True
    
    date_shape = slide.shapes.add_textbox(Inches(12.5), Inches(0.35), Inches(3), Inches(0.5))
    p_date = date_shape.text_frame.paragraphs[0]
    p_date.text = f"Report Date: {date.today().strftime('%Y-%m-%d')}"
    p_date.font.size, p_date.alignment = Pt(14), PP_ALIGN.RIGHT

    # --- KPIs Table ---
    kpi_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(7.5), Inches(0.5))
    kpi_title.text_frame.paragraphs[0].text = "Project Overview"
    kpi_title.text_frame.paragraphs[0].font.bold = True
    
    kpi_table_shape = slide.shapes.add_table(3, 2, Inches(0.5), Inches(1.7), Inches(7.5), Inches(1.5))
    kpi_table = kpi_table_shape.table
    kpi_table.columns[0].width, kpi_table.columns[1].width = Inches(2.5), Inches(5.0)

    set_cell_text(kpi_table.cell(0, 0), "Project Manager", bold=True)
    set_cell_text(kpi_table.cell(0, 1), project_details.get('pm', 'N/A'))
    set_cell_text(kpi_table.cell(1, 0), "Current Phase", bold=True)
    set_cell_text(kpi_table.cell(1, 1), project_details.get('phase', 'N/A'))
    set_cell_text(kpi_table.cell(2, 0), "Overall Health", bold=True)
    set_cell_text(kpi_table.cell(2, 1), project_details.get('health_status', 'N/A'))

    # --- Financials Table ---
    fin_title = slide.shapes.add_textbox(Inches(8.5), Inches(1.2), Inches(7.0), Inches(0.5))
    fin_title.text_frame.paragraphs[0].text = "Financial Summary"
    fin_title.text_frame.paragraphs[0].font.bold = True

    fin_table_shape = slide.shapes.add_table(4, 2, Inches(8.5), Inches(1.7), Inches(7.0), Inches(2.0))
    fin_table = fin_table_shape.table
    fin_table.columns[0].width, fin_table.columns[1].width = Inches(3.0), Inches(4.0)

    set_cell_text(fin_table.cell(0, 0), "Budget (BAC)", bold=True)
    set_cell_text(fin_table.cell(0, 1), f"${project_details.get('budget_usd', 0):,.0f}")
    set_cell_text(fin_table.cell(1, 0), "Actuals (AC)", bold=True)
    set_cell_text(fin_table.cell(1, 1), f"${project_details.get('actuals_usd', 0):,.0f}")
    set_cell_text(fin_table.cell(2, 0), "Predicted EAC (AI)", bold=True)
    set_cell_text(fin_table.cell(2, 1), f"${project_details.get('predicted_eac_usd', 0):,.0f}")
    set_cell_text(fin_table.cell(3, 0), "CPI / SPI", bold=True)
    set_cell_text(fin_table.cell(3, 1), f"{project_details.get('cpi', 0):.2f} / {project_details.get('spi', 0):.2f}")

    # --- Milestones Table ---
    milestone_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(15), Inches(0.5))
    milestone_title.text_frame.paragraphs[0].text = "Upcoming Milestones"
    milestone_title.text_frame.paragraphs[0].font.bold = True

    milestones_to_show = milestones_df.head(5)
    rows = len(milestones_to_show) + 1
    m_table_shape = slide.shapes.add_table(rows, 3, Inches(0.5), Inches(4.0), Inches(15), Inches(0.5 * rows))
    m_table = m_table_shape.table
    m_table.columns[0].width, m_table.columns[1].width, m_table.columns[2].width = Inches(8.0), Inches(4.0), Inches(3.0)
    
    set_cell_text(m_table.cell(0, 0), "Milestone Description", bold=True)
    set_cell_text(m_table.cell(0, 1), "Due Date", bold=True)
    set_cell_text(m_table.cell(0, 2), "Status", bold=True)

    for i, row in enumerate(milestones_to_show.itertuples()):
        set_cell_text(m_table.cell(i + 1, 0), row.milestone)
        set_cell_text(m_table.cell(i + 1, 1), row.due_date)
        set_cell_text(m_table.cell(i + 1, 2), row.status)

    # --- Risks Table ---
    risk_title = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(15), Inches(0.5))
    risk_title.text_frame.paragraphs[0].text = "Top Risks"
    risk_title.text_frame.paragraphs[0].font.bold = True

    risks_to_show = risks_df.head(3)
    if not risks_to_show.empty:
        rows = len(risks_to_show) + 1
        r_table_shape = slide.shapes.add_table(rows, 3, Inches(0.5), Inches(7.0), Inches(15), Inches(0.4 * rows))
        r_table = r_table_shape.table
        r_table.columns[0].width, r_table.columns[1].width, r_table.columns[2].width = Inches(9.0), Inches(3.0), Inches(3.0)
        
        set_cell_text(r_table.cell(0, 0), "Risk Description", bold=True)
        set_cell_text(r_table.cell(0, 1), "Owner", bold=True)
        set_cell_text(r_table.cell(0, 2), "Status", bold=True)

        for i, row in enumerate(risks_to_show.itertuples()):
            set_cell_text(r_table.cell(i + 1, 0), row.description)
            set_cell_text(r_table.cell(i + 1, 1), row.owner)
            set_cell_text(r_table.cell(i + 1, 2), row.status)
    else:
        no_risk_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(15), Inches(0.5))
        no_risk_box.text_frame.paragraphs[0].text = "No open risks logged for this project."
        no_risk_box.text_frame.paragraphs[0].font.italic = True
    
    # --- Save to Buffer ---
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

def generate_board_ready_deck(projects_df: pd.DataFrame, goals_df: pd.DataFrame, alerts: list, plot_utils) -> io.BytesIO:
    """Generates a multi-slide, board-ready executive portfolio summary deck."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(16), Inches(9) # 16:9 aspect ratio

    # --- Slide 1: Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "sPMO Portfolio Review"
    subtitle.text = f"Werfen Autoimmunity Division | {date.today().strftime('%B %Y')}"

    # --- Slide 2: Executive Summary & KPIs ---
    content_slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary & Key Portfolio Indicators"
    
    total_budget = projects_df['budget_usd'].sum()
    active_projects_count = len(projects_df[projects_df['health_status'] != 'Completed'])
    at_risk_count = len(projects_df[projects_df['health_status'] == 'At Risk'])
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(2))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text, p.font.bold, p.font.size = "Portfolio Health at a Glance:", True, Pt(24)
    p = tf.add_paragraph(); p.text = f"Total Active Projects: {active_projects_count}"; p.level = 1
    p = tf.add_paragraph(); p.text = f"Projects At Risk: {at_risk_count}"; p.level = 1
    p = tf.add_paragraph(); p.text = f"Total Portfolio Budget (BAC): ${total_budget:,.0f}"; p.level = 1

    if alerts:
        alert_box = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(14), Inches(3))
        tf = alert_box.text_frame
        p = tf.paragraphs[0]
        p.text, p.font.bold, p.font.size = "Key Automated Alerts:", True, Pt(24)
        for alert in alerts[:3]: # Show top 3 alerts
            p = tf.add_paragraph()
            p.text = f"({alert['type']}) {alert['message']}"; p.level = 1
            if alert['severity'] == 'error': p.font.color.rgb = RGBColor(255, 0, 0)

    # --- Slide 3: Portfolio Landscape Chart ---
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    title.text = "Portfolio Landscape: Strategy vs. Risk"
    fig_bubble = plot_utils.create_portfolio_bubble_chart(projects_df)
    add_plotly_fig_to_slide(fig_bubble, slide, Inches(1.5), Inches(1.5), Inches(13), Inches(7))

    # --- Slide 4: Strategic Alignment ---
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    title.text = "Portfolio Alignment to Strategic Goals"
    
    if not isinstance(goals_df, pd.DataFrame): goals_df = pd.DataFrame(goals_df)
    aligned_df = pd.merge(projects_df, goals_df, left_on='strategic_goal_id', right_on='id', how='left')
    budget_by_goal = aligned_df.groupby('goal')['budget_usd'].sum().reset_index()
    fig_pie = px.pie(budget_by_goal, names='goal', values='budget_usd', title='Portfolio Budget Allocation by Goal')
    add_plotly_fig_to_slide(fig_pie, slide, Inches(4), Inches(1.5), Inches(8), Inches(6))

    # --- Slide 5: High-Risk Projects ---
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    title.text = "Focus Area: High-Risk Projects"
    
    at_risk_df = projects_df[projects_df['health_status'] == 'At Risk'].sort_values('budget_usd', ascending=False)
    if not at_risk_df.empty:
        rows, cols = len(at_risk_df) + 1, 4
        table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(14), Inches(0.5 * rows))
        table = table_shape.table
        
        table.columns[0].width, table.columns[1].width = Inches(5.0), Inches(3.0)
        table.columns[2].width, table.columns[3].width = Inches(3.0), Inches(3.0)
        set_cell_text(table.cell(0,0), 'Project Name', bold=True, font_size=12)
        set_cell_text(table.cell(0,1), 'Project Manager', bold=True, font_size=12)
        set_cell_text(table.cell(0,2), 'Issue Summary (CPI/SPI)', bold=True, font_size=12)
        set_cell_text(table.cell(0,3), 'Budget (USD)', bold=True, font_size=12)
        
        for i, row in enumerate(at_risk_df.itertuples()):
            set_cell_text(table.cell(i+1, 0), row.name)
            set_cell_text(table.cell(i+1, 1), row.pm)
            set_cell_text(table.cell(i+1, 2), f"CPI: {getattr(row, 'cpi', 0):.2f}, SPI: {getattr(row, 'spi', 0):.2f}")
            set_cell_text(table.cell(i+1, 3), f"${row.budget_usd:,.0f}")
    else:
        slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(1)).text_frame.text = "No projects are currently classified as 'At Risk'."

    # --- Save to Buffer ---
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer
